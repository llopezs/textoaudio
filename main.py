import os
import pyttsx3
from docx import Document
import PyPDF2
from pptx import Presentation

def extraer_texto_docx(ruta):
    doc = Document(ruta)
    return "\n".join([p.text for p in doc.paragraphs if p.text.strip() != ""])

def extraer_texto_pdf(ruta):
    texto = ""
    with open(ruta, 'rb') as file:
        reader = PyPDF2.PdfReader(file)
        for page in reader.pages:
            texto += page.extract_text() + "\n"
    return texto

def extraer_texto_pptx(ruta):
    texto = ""
    prs = Presentation(ruta)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texto += shape.text + "\n"
    return texto

def convertir_texto_a_mp3(texto, ruta_salida):
    engine = pyttsx3.init()
    engine.setProperty('rate', 150)
    # Buscar voz en español
    for voz in engine.getProperty('voices'):
        if "spanish" in voz.name.lower():
            engine.setProperty('voice', voz.id)
            break
    engine.save_to_file(texto, ruta_salida)
    engine.runAndWait()

def procesar_archivos(carpeta_entrada, carpeta_salida):
    if not os.path.exists(carpeta_salida):
        os.makedirs(carpeta_salida)

    for archivo in os.listdir(carpeta_entrada):
        ruta_archivo = os.path.join(carpeta_entrada, archivo)
        nombre_base, extension = os.path.splitext(archivo)
        extension = extension.lower()

        texto = ""
        if extension == ".docx":
            texto = extraer_texto_docx(ruta_archivo)
        elif extension == ".pdf":
            texto = extraer_texto_pdf(ruta_archivo)
        elif extension == ".pptx":
            texto = extraer_texto_pptx(ruta_archivo)
        else:
            continue  # Ignorar otros archivos

        if texto.strip():
            ruta_mp3 = os.path.join(carpeta_salida, f"{nombre_base}.mp3")
            print(f"🎤 Generando MP3: {ruta_mp3}")
            convertir_texto_a_mp3(texto, ruta_mp3)

# ===== USO =====
carpeta_entrada = "in"  # Reemplaza con tu ruta
carpeta_salida = "out"    # Reemplaza con tu ruta

procesar_archivos(carpeta_entrada, carpeta_salida)